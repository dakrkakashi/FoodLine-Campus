import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_foodline_pitch_deck(output_path):
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6] # completely blank layout

    # Color Palette
    C_BG = RGBColor(10, 10, 15)          # #0A0A0F
    C_CARD = RGBColor(22, 22, 30)        # #16161E
    C_CARD_BORDER = RGBColor(45, 45, 60) # #2D2D3C
    C_ORANGE = RGBColor(255, 107, 44)    # #FF6B2C
    C_AMBER = RGBColor(255, 179, 71)     # #FFB347
    C_TEAL = RGBColor(0, 212, 170)       # #00D4AA
    C_PURPLE = RGBColor(139, 92, 246)    # #8B5CF6
    C_BLUE = RGBColor(59, 130, 246)      # #3B82F6
    C_RED = RGBColor(233, 69, 96)        # #E94560
    C_TEXT_WHITE = RGBColor(245, 245, 247)
    C_TEXT_GRAY = RGBColor(161, 161, 170)
    C_TEXT_MUTED = RGBColor(113, 113, 122)

    FONT_HEADING = "Outfit"
    FONT_BODY = "Segoe UI"

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, tag_text, title_text, subtitle_text=""):
        # Top Accent Line
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.06))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = C_ORANGE
        top_bar.line.fill.background()

        # Tag
        tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.35))
        tf_tag = tb_tag.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.name = FONT_HEADING
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = C_ORANGE

        # Title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.65))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_HEADING
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = C_TEXT_WHITE

        # Subtitle
        if subtitle_text:
            tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.4))
            tf_sub = tb_sub.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.name = FONT_BODY
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = C_TEXT_GRAY

    def add_card(slide, left, top, width, height, bg_color=C_CARD, border_color=C_CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # =========================================================================
    # SLIDE 1: HERO / TITLE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    # Hero Image on Right Side if available
    img_hero = 'FoodLine_Pitch_Deck/images/hero.jpg'
    if os.path.exists(img_hero):
        slide1.shapes.add_picture(img_hero, Inches(7.5), Inches(1.2), width=Inches(5.0), height=Inches(5.2))
        # Overlay gradient card border
        overlay = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.2), Inches(5.0), Inches(5.2))
        overlay.fill.background()
        overlay.line.color.rgb = C_ORANGE
        overlay.line.width = Pt(2)

    # Tag
    tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(6.5), Inches(0.4))
    p = tb1.text_frame.paragraphs[0]
    p.text = "STARTUP PITCH CONCEPT • CAMPUS DINING 2026–2027"
    p.font.name = FONT_HEADING
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE

    # Brand Title
    tb_b = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.5), Inches(1.1))
    p = tb_b.text_frame.paragraphs[0]
    p.text = "FoodLine"
    p.font.name = FONT_HEADING
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE

    # Tagline
    tb_t = slide1.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(6.5), Inches(0.6))
    p = tb_t.text_frame.paragraphs[0]
    p.text = "Skip the Line, Not the Meal."
    p.font.name = FONT_HEADING
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_WHITE

    # Description
    tb_d = slide1.shapes.add_textbox(Inches(0.8), Inches(3.35), Inches(6.3), Inches(0.9))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    p = tf_d.paragraphs[0]
    p.text = "Next-Generation Campus Pre-Ordering, Smart Slot Throttling & 30-Second Express Counter Collection Ecosystem for Indian Universities."
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = C_TEXT_GRAY

    # 4 Pills
    pills_data = [
        ("📍 Target Pilot Campus", "Sanjivani University (Cafe @7)"),
        ("⚡ Express Handover", "30-Sec Counter QR Pass"),
        ("💰 0% UPI PG Fees", "Direct Merchant UTR Engine"),
        ("📡 Real-Time SSE", "Zero-Refresh Live Tracking")
    ]
    for i, (title, sub) in enumerate(pills_data):
        row = i // 2
        col = i % 2
        px = Inches(0.8 + col * 3.1)
        py = Inches(4.5 + row * 1.1)
        add_card(slide1, px, py, Inches(2.95), Inches(0.95), bg_color=C_CARD)
        tb_p = slide1.shapes.add_textbox(px + Inches(0.12), py + Inches(0.1), Inches(2.7), Inches(0.75))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        p1 = tf_p.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE
        p2 = tf_p.add_paragraph()
        p2.text = sub
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = C_ORANGE if "Target" in title else C_TEXT_GRAY

    # =========================================================================
    # SLIDE 2: PROBLEM STATEMENT
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "Problem Statement", "The Collegiate Break Bottleneck", 
               "Why 4.65 million college students across Maharashtra lose valuable break time daily.")

    cards_p = [
        ("⏳", "40% Break Time Lost", 
         "Short 30-minute campus breaks are crushed by 10–15 minute billing queues before food preparation even begins. Students spend more time waiting than eating.",
         "14.5 min avg wait | 20 min meal window"),
        ("🍳", "Kitchen Rush Chaos", 
         "Canteen operators face uncoordinated demand spikes, leading to kitchen bottlenecks, food wastage, or stockouts — with no visibility into upcoming demand.",
         "Unpredictable spikes | 35% revenue loss"),
        ("💰", "Payment & Cash Friction", 
         "Cash change delays, offline token slips, and high 2% payment gateway fees eat into fragile canteen operating margins on ₹20–50 items.",
         "2% PG fees | Cash bottleneck")
    ]
    for i, (icon, title, desc, badge) in enumerate(cards_p):
        cx = Inches(0.8 + i * 3.95)
        cy = Inches(2.0)
        cw = Inches(3.8)
        ch = Inches(3.6)
        add_card(slide2, cx, cy, cw, ch)
        
        tb = slide2.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.2), cw - Inches(0.4), ch - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = FONT_HEADING
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_TEXT_GRAY

        p3 = tf.add_paragraph()
        p3.space_before = Pt(14)
        p3.text = f"📌 {badge}"
        p3.font.name = FONT_HEADING
        p3.font.size = Pt(10.5)
        p3.font.bold = True
        p3.font.color.rgb = C_ORANGE

    # Bottom Callout Banner
    callout2 = add_card(slide2, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.1), bg_color=RGBColor(25, 20, 25), border_color=C_ORANGE)
    tb_c2 = slide2.shapes.add_textbox(Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.0))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "📊 KEY METRIC"
    p.font.name = FONT_HEADING
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p_b = tf_c2.add_paragraph()
    p_b.text = "Average student spends 14.5 minutes in line for a 20-minute meal window. Result: Students skip meals, arrive late to lectures, and campus food vendors lose over 35% of potential daily turnover."
    p_b.font.name = FONT_BODY
    p_b.font.size = Pt(11.5)
    p_b.font.color.rgb = C_TEXT_WHITE

    # =========================================================================
    # SLIDE 3: PRODUCT SOLUTION
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Product Solution", "The FoodLine Campus Dining Ecosystem", 
               "Four interconnected pillars designed for sub-second collegiate dining efficiency.")

    sol_cards = [
        ("📱", "1. Browse & Order Ahead", "Students browse real-time canteen menus and customize meals directly from classrooms before break bells ring.", "✔ Cloud Sync  ✔ Frictionless"),
        ("🎯", "2. Smart Slot Throttling", "10-minute dynamic break slots balance kitchen load and eliminate rush-hour batching delays.", "✔ Load Balance  ✔ Auto-Slots"),
        ("🔐", "3. 0% Fee UPI QR (Option C)", "Direct vendor UPI scanning + 12-digit UTR validation eliminates payment gateway fees completely.", "✔ Zero Fees  ✔ Direct UTR"),
        ("🚀", "4. 30-Sec Express Handover", "High-contrast dynamic optical QR pass and 4-digit token verification at the dedicated FoodLine counter.", "✔ QR Pass  ✔ 30-Sec Pickup")
    ]
    for i, (icon, title, desc, badges) in enumerate(sol_cards):
        cx = Inches(0.8 + i * 2.95)
        cy = Inches(2.1)
        cw = Inches(2.8)
        ch = Inches(4.6)
        add_card(slide3, cx, cy, cw, ch)

        tb = slide3.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.2), cw - Inches(0.4), ch - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}"
        p.font.size = Pt(28)

        p1 = tf.add_paragraph()
        p1.space_before = Pt(8)
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_GRAY

        p3 = tf.add_paragraph()
        p3.space_before = Pt(14)
        p3.text = badges
        p3.font.name = FONT_HEADING
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = C_TEAL

    # =========================================================================
    # SLIDE 4: FLAGSHIP PILOT PLAN
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Flagship Pilot Plan", "Target Launch: Sanjivani University, Kopargaon", 
               "Proposed deployment model at Cafe @7 designed around actual campus schedule and seed menu.")

    # Left Card: University & Schedule
    add_card(slide4, Inches(0.8), Inches(2.0), Inches(5.7), Inches(4.8))
    tb_p_left = slide4.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.4))
    tf_pl = tb_p_left.text_frame
    tf_pl.word_wrap = True

    p = tf_pl.paragraphs[0]
    p.text = "🏫 Planned Pilot Campus & Outlet"
    p.font.name = FONT_HEADING
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_WHITE

    details = [
        ("University", "Sanjivani University, Kopargaon (MH)"),
        ("Anchor Outlet", "Cafe @7 (Main Campus Cafeteria)"),
        ("Merchant UPI", "Direct Merchant Standee QR (Option C)"),
        ("Break 1 (Lunch)", "11:50 AM – 12:30 PM (40 mins)"),
        ("  • Slot A", "11:50 AM – 12:10 PM"),
        ("  • Slot B", "12:10 PM – 12:30 PM"),
        ("Break 2 (Evening)", "2:30 PM – 2:50 PM (20 mins)")
    ]
    for lbl, val in details:
        p_row = tf_pl.add_paragraph()
        p_row.space_before = Pt(6)
        p_row.text = f"{lbl}: {val}"
        p_row.font.name = FONT_BODY
        p_row.font.size = Pt(12)
        p_row.font.color.rgb = C_ORANGE if "Slot" in lbl or "Break" in lbl else C_TEXT_GRAY

    # Right Card: Seed Menu & Target Pricing
    add_card(slide4, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.8))
    tb_p_right = slide4.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.3), Inches(4.4))
    tf_pr = tb_p_right.text_frame
    tf_pr.word_wrap = True

    p = tf_pr.paragraphs[0]
    p.text = "🍽️ Cafe @7 Verified Menu Highlights"
    p.font.name = FONT_HEADING
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_WHITE

    menu_items = [
        ("Vada Pav / Dabeli / Samosa", "Fast Grab", "₹20"),
        ("Special Campus Tea", "Hot Fuel", "₹10"),
        ("Crispy Masala Dosa", "Bestseller", "₹50"),
        ("Thick Iced Cold Coffee", "Student Fav", "₹50"),
        ("Veg. Cheese Grill Sandwich", "Hot Grill", "₹100"),
        ("Peri Peri French Fries", "Snack Fav", "₹80"),
        ("Veg Fried Momo", "Hot Momo", "₹70"),
        ("KitKat / Oreo Thick Shake", "Shake", "₹90")
    ]
    for name, tag, price in menu_items:
        p_m = tf_pr.add_paragraph()
        p_m.space_before = Pt(5)
        p_m.text = f"• {name}  [{tag}]  —  {price}"
        p_m.font.name = FONT_BODY
        p_m.font.size = Pt(11.5)
        p_m.font.color.rgb = C_TEXT_WHITE

    # =========================================================================
    # SLIDE 5: PAYMENT INNOVATION (OPTION C)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Payment Innovation", "Option C: Zero-Fee UPI & UTR Verification", 
               "Solving the micro-transaction dilemma: 0% payment gateway fees on ₹20 campus items.")

    steps = [
        ("Step 1", "📱", "Scan Standee QR", "Student views Cafe @7 UPI QR Standee image and UPI ID displayed in the FoodLine app."),
        ("Step 2", "⚡", "Pay via Any UPI App", "Pays exact total directly via GPay, PhonePe, Paytm, or CRED with zero platform fees."),
        ("Step 3", "🔑", "Submit 12-Digit UTR", "Inputs bank UTR reference number into FoodLine verification modal for real-time validation."),
        ("Step 4", "🎟️", "Generate Token & Pass", "Backend validates uniqueness, prevents replay, and issues verified digital token #FL-XXXX.")
    ]
    for i, (st, icon, title, desc) in enumerate(steps):
        cx = Inches(0.8 + i * 2.95)
        cy = Inches(2.0)
        cw = Inches(2.8)
        ch = Inches(3.6)
        add_card(slide5, cx, cy, cw, ch)

        tb = slide5.shapes.add_textbox(cx + Inches(0.18), cy + Inches(0.18), cw - Inches(0.36), ch - Inches(0.36))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{st} • {icon}"
        p.font.name = FONT_HEADING
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_TEAL

        p1 = tf.add_paragraph()
        p1.space_before = Pt(8)
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_GRAY

    # Callout Banner
    add_card(slide5, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.1), bg_color=RGBColor(20, 30, 25), border_color=C_TEAL)
    tb_c5 = slide5.shapes.add_textbox(Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.0))
    tf_c5 = tb_c5.text_frame
    tf_c5.word_wrap = True
    p = tf_c5.paragraphs[0]
    p.text = "🏆 WHY OPTION C WINS ON CAMPUS"
    p.font.name = FONT_HEADING
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_TEAL
    p_b = tf_c5.add_paragraph()
    p_b.text = "Traditional Razorpay/Paytm PG charges 2.0% + GST per transaction — unfeasible for low ticket ₹20–50 canteen items. Option C ensures 100% merchant payout directly to the canteen's bank account with zero intermediary commissions."
    p_b.font.name = FONT_BODY
    p_b.font.size = Pt(11.5)
    p_b.font.color.rgb = C_TEXT_WHITE

    # =========================================================================
    # SLIDE 6: KITCHEN INTELLIGENCE (KDS & SSE)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "Kitchen Intelligence", "Kitchen Display System (KDS) & Real-Time Sync", 
               "Empowering cafeteria cooks with automated batch forecasting and real-time handovers.")

    kds_features = [
        ("📈", "Smart Slot Batching", "Groups incoming orders by upcoming 10-minute break slots so kitchen staff can batch-cook (e.g., 15 Dosas for 11:50 AM). Reduces waste and prep chaos."),
        ("📡", "Server-Sent Events (SSE)", "Real-time bi-directional pipeline updates order status from PREPARING → READY with zero screen refreshes. Students see live progress."),
        ("📷", "Optical QR Handover", "Express counter staff scan student digital QR code or enter 4-digit PIN for sub-30-second meal handovers. No paper tokens needed."),
        ("📦", "Instant Inventory Toggle", "1-tap out-of-stock switcher prevents students from ordering sold-out dishes during sudden surges. Real-time menu sync across all devices.")
    ]
    for i, (icon, title, desc) in enumerate(kds_features):
        row = i // 2
        col = i % 2
        cx = Inches(0.8 + col * 5.95)
        cy = Inches(2.1 + row * 2.4)
        add_card(slide6, cx, cy, Inches(5.75), Inches(2.2))

        tb = slide6.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.2), Inches(5.35), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = FONT_HEADING
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_WHITE

        p1 = tf.add_paragraph()
        p1.space_before = Pt(8)
        p1.text = desc
        p1.font.name = FONT_BODY
        p1.font.size = Pt(12)
        p1.font.color.rgb = C_TEXT_GRAY

    # =========================================================================
    # SLIDE 7: MARKET OPPORTUNITY (TAM / SAM / SOM)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7)
    add_header(slide7, "Market Opportunity", "Collegiate Foodservice Market in India", 
               "Massive unorganized campus dining market ready for digital transformation.")

    m_cards = [
        ("₹18,000+ Cr", "TAM — Total Addressable Market", 
         "45,000+ colleges across India with on-campus cafeterias and dining halls serving over 40 million students annually.", 
         "$2.2 Billion Market | 45K+ Colleges"),
        ("₹2,400 Cr", "SAM — Serviceable Addressable", 
         "5,380 colleges & universities in Maharashtra with 46.5 Lakh (4.65M) enrolled collegiate students.", 
         "$290 Million Market | 5,380 Colleges"),
        ("₹35 Cr", "SOM — Target Launch Reach", 
         "Targeting 50 university campuses across Pune, Kopargaon, Ahmednagar, Nashik, and Mumbai university hubs in Year 1-2.", 
         "$4.2M GMV Target | 50 Campuses")
    ]
    for i, (val, title, desc, badge) in enumerate(m_cards):
        cx = Inches(0.8 + i * 3.95)
        cy = Inches(2.1)
        cw = Inches(3.8)
        ch = Inches(4.8)
        add_card(slide7, cx, cy, cw, ch)

        tb = slide7.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.3), cw - Inches(0.4), ch - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = val
        p.font.name = FONT_HEADING
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = C_ORANGE

        p1 = tf.add_paragraph()
        p1.space_before = Pt(10)
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_TEXT_GRAY

        p3 = tf.add_paragraph()
        p3.space_before = Pt(16)
        p3.text = f"🚀 {badge}"
        p3.font.name = FONT_HEADING
        p3.font.size = Pt(10.5)
        p3.font.bold = True
        p3.font.color.rgb = C_AMBER

    # =========================================================================
    # SLIDE 8: BUSINESS MODEL
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8)
    add_header(slide8, "Business Model", "Multi-Stream Sustainable Revenue Model", 
               "Zero barrier to vendor entry paired with scalable transaction and SaaS revenue.")

    bm_streams = [
        ("🆓", "1. Free Canteen Onboarding", "Zero setup fees and zero hardware requirement guarantees 100% canteen vendor adoption across campuses.", "₹0 Setup  •  No Hardware Barrier"),
        ("📊", "2. Pickup Commission (3-5%)", "Small platform fee per successful order covering hosting, SSE infrastructure, and continuous support.", "3-5% Per Order  •  Volume Scaled"),
        ("🏃", "3. Campus Delivery Fee", "Optional in-campus delivery (₹10–20) to library and hostel drop zones, powered by student runners.", "₹10-20 Fee  •  Peer Logistics"),
        ("💎", "4. Merchant Analytics SaaS", "Predictive demand forecasting, wastage reduction, and menu engineering dashboards for vendors.", "Premium SaaS  •  AI Reports")
    ]
    for i, (icon, title, desc, badges) in enumerate(bm_streams):
        cx = Inches(0.8 + i * 2.95)
        cy = Inches(2.1)
        cw = Inches(2.8)
        ch = Inches(4.6)
        add_card(slide8, cx, cy, cw, ch)

        tb = slide8.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.2), cw - Inches(0.4), ch - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}"
        p.font.size = Pt(28)

        p1 = tf.add_paragraph()
        p1.space_before = Pt(8)
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_GRAY

        p3 = tf.add_paragraph()
        p3.space_before = Pt(14)
        p3.text = badges
        p3.font.name = FONT_HEADING
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = C_TEAL

    # =========================================================================
    # SLIDE 9: COMPETITIVE MATRIX
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9)
    add_header(slide9, "Competitive Matrix", "Why Existing Delivery Apps Fail on Campus", 
               "FoodLine is architected specifically for collegiate micro-breaks and low-AOV dynamics.")

    # Table Shape
    rows = 8
    cols = 4
    left = Inches(0.8)
    top = Inches(2.0)
    width = Inches(11.7)
    height = Inches(4.8)

    table_shape = slide9.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(2.6)
    table.columns[2].width = Inches(2.8)
    table.columns[3].width = Inches(3.1)

    headers = ["Feature / Attribute", "Traditional Canteen", "Swiggy / Zomato", "FoodLine Campus ✓"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 30, 45)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = FONT_HEADING
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_ORANGE if j == 3 else C_TEXT_WHITE

    matrix_rows = [
        ("Break Window Alignment", "❌ 15-min queue chaos", "❌ 35-45 min delivery", "✅ 10-min smart slot pickup"),
        ("Average Order Value", "₹20 – ₹60", "❌ ₹200+ (High minimums)", "✅ No minimum (₹15 Chai)"),
        ("Delivery / Platform Fees", "Zero (offline only)", "❌ ₹30-60 + 25% take", "✅ 0% Fee UPI QR"),
        ("Kitchen Load Throttling", "❌ Overload & delays", "❌ Uncoordinated spikes", "✅ Algorithmic slot throttle"),
        ("Campus Integration", "❌ No digital layer", "❌ External delivery only", "✅ Native campus system"),
        ("Real-Time Order Tracking", "❌ No visibility", "✅ Delivery GPS", "✅ SSE live status"),
        ("Vendor Onboarding Cost", "N/A", "❌ Commission + listing", "✅ ₹0 Free onboarding")
    ]
    for i, row in enumerate(matrix_rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(18, 18, 25) if (i % 2 == 0) else RGBColor(24, 24, 34)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            p.font.color.rgb = C_TEAL if j == 3 else C_TEXT_WHITE

    # =========================================================================
    # SLIDE 10: TARGET IMPACT & KPIS
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10)
    add_header(slide10, "Impact & Projected Benchmarks", "Target Operational Breakthroughs", 
               "Quantifiable performance benchmarks engineered for campus cafeterias.")

    kpis = [
        ("≥ 12 Min", "Target Time Saved", "Students gain back 40%+ of their break time to socialize, study, or relax instead of standing in queues."),
        ("2.5×", "Peak Capacity Boost", "Canteens fulfill 2.5× more orders in the same 30-minute rush window without expanding space or staff."),
        ("≤ 30 Sec", "Express Pickup Goal", "Instant QR scan & 4-digit PIN verification ensures rapid pickup with zero billing delays at the counter."),
        ("≥ 65%", "Projected Retention", "Daily campus dining habits create industry-leading organic retention rates — no marketing spend needed.")
    ]
    for i, (val, lbl, desc) in enumerate(kpis):
        cx = Inches(0.8 + i * 2.95)
        cy = Inches(2.2)
        cw = Inches(2.8)
        ch = Inches(4.4)
        add_card(slide10, cx, cy, cw, ch)

        tb = slide10.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.3), cw - Inches(0.4), ch - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = val
        p.font.name = FONT_HEADING
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = C_TEAL

        p1 = tf.add_paragraph()
        p1.space_before = Pt(12)
        p1.text = lbl
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = C_TEXT_GRAY

    # =========================================================================
    # SLIDE 11: ENGINEERING ARCHITECTURE (TECH STACK)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide11)
    add_header(slide11, "Engineering Architecture", "Modern Enterprise Technology Stack", 
               "Engineered for high concurrency, sub-second latency, and responsive mobile experiences.")

    tech_blocks = [
        ("⚙️", "Frontend Engine", "Production-grade frontend with server-side rendering, strict type safety, and spring-physics animations for a native-app feel.", "Next.js 15+  •  React 19  •  TypeScript  •  Tailwind v4"),
        ("🧩", "Component System", "Accessible, themeable component library with high-contrast OKLCH color tokens and custom UI primitives.", "shadcn/ui  •  Radix UI  •  OKLCH Tokens  •  Nova"),
        ("🔌", "Backend & Real-Time", "Serverless API handlers with persistent SSE connections for zero-latency bi-directional order updates.", "API Routes  •  SSE Pipeline  •  Supabase  •  PostgreSQL"),
        ("🧠", "Knowledge Graph", "Graphify AST mapping for automated code intelligence, clean architecture enforcement, and dev velocity.", "Graphify  •  AST Analysis  •  Code Intelligence")
    ]
    for i, (icon, title, desc, badges) in enumerate(tech_blocks):
        row = i // 2
        col = i % 2
        cx = Inches(0.8 + col * 5.95)
        cy = Inches(2.1 + row * 2.4)
        add_card(slide11, cx, cy, Inches(5.75), Inches(2.2))

        tb = slide11.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.18), Inches(5.35), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = FONT_HEADING
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_WHITE

        p1 = tf.add_paragraph()
        p1.space_before = Pt(6)
        p1.text = desc
        p1.font.name = FONT_BODY
        p1.font.size = Pt(11.5)
        p1.font.color.rgb = C_TEXT_GRAY

        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        p2.text = f"⚡ {badges}"
        p2.font.name = FONT_HEADING
        p2.font.size = Pt(10.5)
        p2.font.bold = True
        p2.font.color.rgb = C_BLUE

    # =========================================================================
    # SLIDE 12: ROADMAP & TIMELINE
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide12)
    add_header(slide12, "Roadmap & Timeline", "From Concept to Pan-India Scale", 
               "A 3-phase execution strategy designed for rapid, capital-efficient campus expansion.")

    phases = [
        ("Phase 1", "Pilot Launch & Validation", "Target: Sanjivani University (Cafe @7)", 
         ["Launch flagship pilot at Sanjivani University", "Target 500+ daily active student orders", "Validate Option C UPI flow with campus vendors", "Refine KDS and SSE infrastructure in real conditions", "Onboard additional campus canteens"]),
        ("Phase 2", "City-Cluster Expansion", "Pune, Nashik & Mumbai Clusters", 
         ["Expand across regional university hubs", "Target 50 campus deployments", "Launch campus delivery with student runners", "Introduce Merchant Analytics SaaS", "Target ₹35 Cr annualized GMV"]),
        ("Phase 3", "National Standardization", "5,000+ Colleges Across India", 
         ["Scale across major universities nationwide", "Partner with university administrations", "Launch FoodLine API for canteen POS systems", "AI-powered demand prediction engine", "Target ₹500 Cr+ annual GMV"])
    ]
    for i, (ph, title, sub, bullets) in enumerate(phases):
        cx = Inches(0.8 + i * 3.95)
        cy = Inches(2.1)
        cw = Inches(3.8)
        ch = Inches(4.8)
        add_card(slide12, cx, cy, cw, ch)

        tb = slide12.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.25), cw - Inches(0.4), ch - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = ph.upper()
        p.font.name = FONT_HEADING
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_ORANGE

        p1 = tf.add_paragraph()
        p1.space_before = Pt(4)
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        p2.text = sub
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_AMBER

        for b in bullets:
            pb = tf.add_paragraph()
            pb.space_before = Pt(6)
            pb.text = f"• {b}"
            pb.font.name = FONT_BODY
            pb.font.size = Pt(10.5)
            pb.font.color.rgb = C_TEXT_GRAY

    # =========================================================================
    # SLIDE 13: FINANCIAL PROJECTIONS
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide13)
    add_header(slide13, "Financial Projections", "Revenue & Growth Trajectory", 
               "Projections based on target unit economics and Maharashtra campus density.")

    # 3 Revenue Cards
    fin_years = [
        ("Year 1 (Pilot Phase)", "₹1.2 Cr", "Gross Target Revenue", "Pilot Target"),
        ("Year 2 (Expansion)", "₹8.5 Cr", "Gross Target Revenue", "↑ 608% YoY Growth"),
        ("Year 3 (Scale)", "₹35 Cr", "Gross Target Revenue", "↑ 312% YoY Growth")
    ]
    for i, (yr, amt, lbl, grw) in enumerate(fin_years):
        cx = Inches(0.8 + i * 3.95)
        cy = Inches(2.0)
        cw = Inches(3.8)
        ch = Inches(1.8)
        add_card(slide13, cx, cy, cw, ch)

        tb = slide13.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4), ch - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = yr
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_AMBER

        p1 = tf.add_paragraph()
        p1.space_before = Pt(4)
        p1.text = amt
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.text = f"{lbl}  [{grw}]"
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = C_TEAL

    # 3 Unit Economics Cards
    unit_data = [
        ("Avg Order Value", "₹45", "Per Transaction"),
        ("Target Daily Orders", "12,000+", "Across 50 Campuses (Y2)"),
        ("Target Gross Margin", "72%", "Asset-Light Model")
    ]
    for i, (title, amt, sub) in enumerate(unit_data):
        cx = Inches(0.8 + i * 3.95)
        cy = Inches(4.0)
        cw = Inches(3.8)
        ch = Inches(1.6)
        add_card(slide13, cx, cy, cw, ch)

        tb = slide13.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4), ch - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_GRAY

        p1 = tf.add_paragraph()
        p1.space_before = Pt(4)
        p1.text = amt
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = C_TEXT_MUTED

    # Bottom Unit Economics Callout
    add_card(slide13, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.1), bg_color=RGBColor(25, 20, 20), border_color=C_ORANGE)
    tb_c13 = slide13.shapes.add_textbox(Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.0))
    tf_c13 = tb_c13.text_frame
    tf_c13.word_wrap = True
    p = tf_c13.paragraphs[0]
    p.text = "💡 UNIT ECONOMICS ADVANTAGE"
    p.font.name = FONT_HEADING
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p_b = tf_c13.add_paragraph()
    p_b.text = "CAC: ₹0 (organic campus adoption via QR standees) • LTV: ₹2,400+/year (avg student places 200+ orders/year) • LTV:CAC Ratio: ∞. Break-even expected within 8 months of campus deployment."
    p_b.font.name = FONT_BODY
    p_b.font.size = Pt(11.5)
    p_b.font.color.rgb = C_TEXT_WHITE

    # =========================================================================
    # SLIDE 14: THE TEAM
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide14)
    add_header(slide14, "The Team", "Built by Students, for Students", 
               "A founding team with deep campus insight and full-stack engineering capabilities.")

    team = [
        ("SN", "Shivam Nirmal", "Founder & CEO", 
         "Full-stack engineer and campus operations strategist. Architect of FoodLine's Option C payment engine and SSE infrastructure."),
        ("TC", "Technical Co-Founder", "Chief Technology Officer (CTO)", 
         "Systems architect specializing in real-time event-driven architectures, database optimization, and mobile-first experiences."),
        ("OL", "Operations Lead", "Chief Operating Officer (COO)", 
         "Campus engagement expert managing vendor relationships, student ambassador programs, and ground-level pilot execution.")
    ]
    for i, (avatar, name, role, desc) in enumerate(team):
        cx = Inches(0.8 + i * 3.95)
        cy = Inches(2.1)
        cw = Inches(3.8)
        ch = Inches(4.8)
        add_card(slide14, cx, cy, cw, ch)

        # Avatar Circle
        av = slide14.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.3), cy + Inches(0.4), Inches(0.9), Inches(0.9))
        av.fill.solid()
        av.fill.fore_color.rgb = C_ORANGE if i == 0 else (C_TEAL if i == 1 else C_PURPLE)
        av.line.fill.background()
        p_av = av.text_frame.paragraphs[0]
        p_av.text = avatar
        p_av.alignment = PP_ALIGN.CENTER
        p_av.font.name = FONT_HEADING
        p_av.font.size = Pt(16)
        p_av.font.bold = True
        p_av.font.color.rgb = C_BG

        tb = slide14.shapes.add_textbox(cx + Inches(0.3), cy + Inches(1.5), cw - Inches(0.6), ch - Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_HEADING
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_WHITE

        p1 = tf.add_paragraph()
        p1.space_before = Pt(2)
        p1.text = role
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = C_ORANGE if i == 0 else (C_TEAL if i == 1 else C_PURPLE)

        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = C_TEXT_GRAY

    # =========================================================================
    # SLIDE 15: CLOSING & CALL TO ACTION
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide15)

    # Center Hero Content
    tb_c_tag = slide15.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.4))
    p = tb_c_tag.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "THE FUTURE OF CAMPUS DINING"
    p.font.name = FONT_HEADING
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE

    tb_c_title = slide15.shapes.add_textbox(Inches(1.5), Inches(1.4), Inches(10.3), Inches(1.2))
    p = tb_c_title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "The Future of Campus Dining Starts Here"
    p.font.name = FONT_HEADING
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_WHITE

    tb_c_sub = slide15.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(10.3), Inches(0.6))
    p = tb_c_sub.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "A Startup Concept Ready to Scale: From Sanjivani University to 5,000+ Campuses Across India."
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.color.rgb = C_TEXT_GRAY

    # 3 Summary Cards
    c_steps = [
        ("Phase 1", "Flagship Pilot", "Sanjivani University (Cafe @7)"),
        ("Phase 2", "City Expansion", "50 Campuses across MH"),
        ("Phase 3", "National Scale", "5,000+ Colleges Pan-India")
    ]
    for i, (ph, t, sub) in enumerate(c_steps):
        cx = Inches(2.2 + i * 3.1)
        cy = Inches(3.4)
        add_card(slide15, cx, cy, Inches(2.8), Inches(1.3))
        tb = slide15.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.15), Inches(2.5), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = ph.upper()
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_ORANGE
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.text = t
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = sub
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10)
        p2.font.color.rgb = C_TEXT_MUTED

    # Bottom Branding & CTA
    tb_bot = slide15.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.8))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    p = tf_bot.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "🚀 FoodLine — Skip the Line, Not the Meal."
    p.font.name = FONT_HEADING
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_WHITE

    p_sub = tf_bot.add_paragraph()
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.text = "Prototype Architecture Ready • Next.js 15 + React 19 + Supabase"
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = C_TEXT_MUTED

    # CTA Button
    btn = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.15), Inches(6.0), Inches(3.0), Inches(0.7))
    btn.fill.solid()
    btn.fill.fore_color.rgb = C_ORANGE
    btn.line.fill.background()
    p_btn = btn.text_frame.paragraphs[0]
    p_btn.alignment = PP_ALIGN.CENTER
    p_btn.text = "Let's Talk  →"
    p_btn.font.name = FONT_HEADING
    p_btn.font.size = Pt(14)
    p_btn.font.bold = True
    p_btn.font.color.rgb = C_BG

    # Save presentation
    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation: {output_path}")

if __name__ == '__main__':
    out_file = 'FoodLine_Master_Pitch_Deck_2026.pptx'
    create_foodline_pitch_deck(out_file)
